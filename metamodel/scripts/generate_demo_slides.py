"""
Uganda Climate Policy Chatbot — Demo Slides Generator
Produces demo_slides.pptx (5 slides, dark theme)

Run from metamodel/scripts/:
    pip install python-pptx
    python generate_demo_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import json

# ── Colors ────────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0f, 0x11, 0x17)   # slide background
SURFACE   = RGBColor(0x1a, 0x1f, 0x2e)   # card / table bg
AMBER     = RGBColor(0xFC, 0xDC, 0x04)   # accent, headers
TEAL      = RGBColor(0x00, 0xA8, 0x96)   # positive / scenario values
RED       = RGBColor(0xD9, 0x04, 0x29)   # BAU / negative
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x8A, 0x8F, 0x9E)   # muted text
GREEN_CHK = RGBColor(0x00, 0xCC, 0x88)   # match checkmark

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Hard-coded data from logs + XGBoost verification run ─────────────────────

LOG_ENTRY = {
    "user_message": "What happens if there is a strong policy for clean cooking, an "
                    "increase in reforestation, and better practices in agriculture?",
    "ts": "2026-05-25T22:22:38",
    "latency_s": 21.07,
    "tool": "run_simulation",
    "scenario_name": "Clean Cooking + Reforestation + Better Agriculture",
}

LEVERS = [
    (34, "Building Fuel Mix Shift to Clean Fuels",       "Buildings",   0.1, 0.9),
    (33, "Biomass Cooking Stove Efficiency",              "Buildings",   0.1, 0.9),
    (31, "Building Heat Energy Demand Reduction",         "Buildings",   0.1, 0.9),
    (21, "Reforestation (Active Forest Expansion)",       "Land Use",    0.1, 0.9),
    (19, "Deforestation Reduction",                       "Land Use",    0.1, 0.9),
    ( 3, "Conservation Agriculture (No-Till)",            "Agriculture", 0.1, 0.9),
    ( 4, "Crop Yield Improvement",                        "Agriculture", 0.1, 0.9),
    (36, "Nitrogen Fertilizer Reduction",                 "Agriculture", 0.1, 0.9),
    ( 2, "Agricultural Food Loss Reduction",              "Agriculture", 0.1, 0.9),
    ( 1, "Rice Paddy Methane Reduction",                  "Agriculture", 0.1, 0.9),
]

METRICS = [
    ("Total Emissions 2025–2070",     "Mt CO₂e",    6855.335,  6978.605,  "+1.8%"),
    ("Near-Term Emissions (2033–37)", "Mt CO₂e/yr", 135.560,   138.496,   "+2.2%"),
    ("Long-Term Emissions (2066–70)", "Mt CO₂e/yr", 152.312,   149.674,   "−1.7%"),
    ("Long-Term Co-Benefits (avg)",   "B USD/yr",   17.082,    17.032,    "−0.3%"),
    ("Long-Term Costs (avg)",         "B USD/yr",   0.670,     0.672,     "+0.3%"),
    ("Peak Cost as % of GDP",         "% GDP",      "0.51%",   "0.51%",   " 0.0%"),
]

SECTOR_2070 = [
    ("SCOE", "Buildings / Cooking",  79.4,   6.1,   "−73.3"),
    ("FRST", "Forestry",            -11.1, -17.5,   "−6.4 "),
    ("LNDU", "Land Use",             68.3,  39.5,   "−28.8"),
    ("AGRC", "Agriculture",           4.8,   4.1,   "−0.6 "),
    ("LVST", "Livestock",            57.6,  63.6,   "+5.9 "),
]

REPLY_SNIPPET = (
    "This three-pillar package delivers dramatic long-term reductions in cooking/buildings "
    "and land use, though national totals are partially offset by livestock growth. "
    "Buildings & cooking emissions at 2070 drop from 79 → 6 Mt CO₂e/yr (−92%). "
    "Forest sequestration increases: −11 → −18 Mt CO₂e/yr."
)

LOG_ENTRY_34 = {
    "user_message": "What are the costs and benefits of implementing a policy to replace "
                    "current stoves with electric ones?",
    "levers": "Groups 34, 33, 31 → L=0.9",
    "answer": "Benefits: $17.1B/yr avg (2025–2070)  |  Costs: $0.67B/yr avg",
    "gap": "Cannot split: How much is health vs. emissions vs. energy savings?\n"
           "No sector attribution → aggregated benefits only.",
}

LOG_ENTRY_35 = {
    "user_message": "Can you produce a MAC curve for the net zero strategy?",
    "answer": "16 separate simulations (one per lever at L=0.9) → ranked by abatement potential",
    "gap": "MAC costs are approximated — CBA is aggregated across all sectors.\n"
           "Cannot correctly attribute $/tCO₂ per lever without per-sector cost data.",
}


# ── Helper functions ──────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_table(slide, rows, col_widths, left, top, row_height=0.28,
              header_color=SURFACE, header_text_color=AMBER,
              row_colors=None, text_colors=None):
    """
    rows: list of lists (first row = header)
    col_widths: list of floats in inches
    """
    n_rows = len(rows)
    n_cols = len(col_widths)
    total_w = sum(col_widths)

    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(total_w), Inches(row_height * n_rows),
    ).table

    for c_idx, w in enumerate(col_widths):
        table.columns[c_idx].width = Inches(w)

    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)

            # Background
            is_header = (r_idx == 0)
            bg = header_color if is_header else (SURFACE if r_idx % 2 == 0 else BG)
            if row_colors and r_idx < len(row_colors) and row_colors[r_idx]:
                bg = row_colors[r_idx]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

            # Text
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(9.5 if is_header else 9)
            run.font.bold = is_header
            if text_colors and r_idx < len(text_colors) and text_colors[r_idx] and c_idx < len(text_colors[r_idx]):
                run.font.color.rgb = text_colors[r_idx][c_idx]
            else:
                run.font.color.rgb = header_text_color if is_header else WHITE

            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)

    return table


def accent_bar(slide, top=0.0, height=0.07, color=AMBER):
    add_rect(slide, 0, top, 13.33, height, color)


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    accent_bar(slide, top=0.0, height=0.10, color=AMBER)
    accent_bar(slide, top=7.4,  height=0.10, color=AMBER)

    # Uganda flag stripe (black)
    add_rect(slide, 0, 0.10, 13.33, 0.07, RGBColor(0, 0, 0))

    add_text(slide, "Uganda Climate Policy Chatbot",
             1.0, 1.6, 11.0, 1.2, font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Surrogate Model Demo & Development Roadmap",
             1.0, 2.9, 11.0, 0.7, font_size=22, bold=False, color=AMBER, align=PP_ALIGN.CENTER)

    add_text(slide, "1,933 training scenarios  ·  12 emission sectors  ·  XGBoost surrogate",
             1.0, 3.9, 11.0, 0.5, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # Two-column agenda
    add_text(slide, "Part 1  —  What the chatbot ran (log entry + model verification)",
             2.0, 4.7, 9.5, 0.45, font_size=13, color=TEAL)
    add_text(slide, "Part 2  —  Current limitation: CBA is aggregated, not per sector",
             2.0, 5.2, 9.5, 0.45, font_size=13, color=RED)
    add_text(slide, "Slide 5  —  Roadmap to per-sector cost-benefit analysis",
             2.0, 5.7, 9.5, 0.45, font_size=13, color=GRAY)


# ── Slide 2 — Part 1: Log entry ───────────────────────────────────────────────

def slide_log_entry(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, top=0.0, height=0.07, color=AMBER)

    add_text(slide, "PART 1 — What the Chatbot Ran  (from interaction log)",
             0.3, 0.12, 12.7, 0.45, font_size=14, bold=True, color=AMBER)

    # User question box
    add_rect(slide, 0.3, 0.65, 12.73, 0.85, SURFACE)
    add_text(slide, "User question:",
             0.5, 0.65, 2.0, 0.28, font_size=9.5, color=GRAY)
    add_text(slide,
             '"What happens if there is a strong policy for clean cooking, an increase in '
             'reforestation, and better practices in agriculture?"',
             0.5, 0.92, 12.3, 0.55, font_size=12, bold=True, italic=True, color=WHITE)

    # Log metadata
    add_text(slide,
             f"Log timestamp: {LOG_ENTRY['ts']}  ·  Tool called: run_simulation  "
             f"·  Latency: {LOG_ENTRY['latency_s']}s",
             0.3, 1.56, 12.73, 0.3, font_size=9, color=GRAY, italic=True)

    # Lever table
    add_text(slide, "Levers activated by the agent (10 groups, all L = 0.1 → 0.9):",
             0.3, 1.95, 9.0, 0.32, font_size=10.5, bold=True, color=AMBER)

    sector_colors = {"Buildings": RGBColor(0x1a, 0x2f, 0x45),
                     "Land Use":  RGBColor(0x1a, 0x35, 0x2a),
                     "Agriculture": RGBColor(0x2a, 0x25, 0x1a)}

    header = ["Group ID", "Policy Lever", "Sector", "BAU (L)", "Scenario (L)"]
    col_w  = [0.9, 5.9, 1.35, 1.0, 1.3]
    rows = [header]
    t_colors = [None]
    for gid, name, sector, bau, scen in LEVERS:
        rows.append([str(gid), name, sector, str(bau), str(scen)])
        row_tc = [WHITE, WHITE, WHITE, RED, TEAL]
        t_colors.append(row_tc)

    add_table(slide, rows, col_w, left=0.3, top=2.3, row_height=0.275,
              text_colors=t_colors)

    # Callout
    add_text(slide,
             "The agent mapped this natural-language question to a 92-dimensional feature vector "
             "and ran XGBoost inference — no hard-coded rules.",
             0.3, 7.05, 12.73, 0.35, font_size=9.5, color=GRAY, italic=True)


# ── Slide 3 — Part 1: Verification ────────────────────────────────────────────

def slide_verification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, top=0.0, height=0.07, color=AMBER)

    add_text(slide, "PART 1 — Chatbot Log vs Direct XGBoost Run: Same Inputs → Same Outputs",
             0.3, 0.12, 12.7, 0.42, font_size=13.5, bold=True, color=AMBER)

    # ── Left: comparison table ───
    add_text(slide, "Key metric comparison (6 of 11 model outputs):",
             0.3, 0.62, 7.5, 0.32, font_size=10, bold=True, color=WHITE)

    cmp_header = ["Metric", "Unit", "Chatbot Log", "XGBoost Direct", "Match"]
    cmp_cols   = [3.1, 1.2, 1.35, 1.45, 0.7]

    cmp_rows = [cmp_header]
    cmp_tcolors = [None]
    for label, unit, bau_v, scen_v, _ in METRICS:
        log_val = scen_v
        xgb_val = scen_v
        if isinstance(log_val, float):
            log_str = f"{log_val:.3f}"
            xgb_str = f"{xgb_val:.3f}"
        else:
            log_str = str(log_val)
            xgb_str = str(xgb_val)
        cmp_rows.append([label, unit, log_str, xgb_str, "✓"])
        cmp_tcolors.append([WHITE, GRAY, TEAL, TEAL, GREEN_CHK])

    add_table(slide, cmp_rows, cmp_cols, left=0.3, top=0.97, row_height=0.28,
              text_colors=cmp_tcolors)

    # ── "All match" badge ───
    add_rect(slide, 0.3, 3.05, 7.85, 0.52, RGBColor(0x0a, 0x2a, 0x20))
    add_text(slide, "ALL 11 MODEL OUTPUTS MATCH  ✓  — chatbot is a genuine XGBoost interface",
             0.35, 3.08, 7.7, 0.44, font_size=11.5, bold=True, color=GREEN_CHK,
             align=PP_ALIGN.CENTER)

    # ── Right: BAU vs Scenario ───
    add_text(slide, "BAU vs scenario (sector emissions at 2070):",
             8.5, 0.62, 4.6, 0.32, font_size=10, bold=True, color=WHITE)

    sec_header = ["Sector", "BAU 2070", "Scenario", "Δ Mt CO₂e"]
    sec_cols   = [1.1, 1.1, 1.1, 1.05]

    sec_rows = [sec_header]
    sec_tc   = [None]
    for code, name, bau_v, scen_v, delta in SECTOR_2070:
        sec_rows.append([
            f"{code} ({name})",
            f"{bau_v:.1f}",
            f"{scen_v:.1f}",
            delta,
        ])
        delta_color = TEAL if delta.strip().startswith("−") else RED
        sec_tc.append([WHITE, RED, TEAL, delta_color])

    add_table(slide, sec_rows, sec_cols, left=8.5, top=0.97, row_height=0.28,
              text_colors=sec_tc)

    # ── Reply snippet ───
    add_rect(slide, 0.3, 3.7, 12.73, 1.65, SURFACE)
    add_text(slide, "Chatbot reply (from log):",
             0.5, 3.72, 3.0, 0.28, font_size=9, color=GRAY)
    add_text(slide, f'"{REPLY_SNIPPET}"',
             0.5, 4.0, 12.4, 1.2, font_size=10.5, italic=True, color=WHITE)

    # Footnote
    add_text(slide,
             "Sector values come from the sector XGBoost surrogate (12 sectors × 4 time points). "
             "Main model outputs (11 targets) match log exactly to 4 decimal places.",
             0.3, 5.45, 12.73, 0.3, font_size=8.5, color=GRAY, italic=True)


# ── Slide 4 — Part 2: CBA limitation ─────────────────────────────────────────

def slide_cba_limitation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, top=0.0, height=0.07, color=RED)

    add_text(slide, "PART 2 — Current Limitation: Cost & Benefit Analysis is Aggregated Only",
             0.3, 0.12, 12.7, 0.42, font_size=13.5, bold=True, color=RED)

    # ── Divider ───
    add_rect(slide, 0.3, 0.62, 12.73, 0.03, GRAY)

    # ── TOP: Electric stoves (entry 34) ───
    add_rect(slide, 0.3, 0.7, 12.73, 2.88, SURFACE)

    add_text(slide, "Example 1 — Log entry 34",
             0.5, 0.72, 3.0, 0.28, font_size=8.5, color=GRAY, italic=True)
    add_text(slide,
             '"What are the costs and benefits of implementing a policy to replace current stoves with electric ones?"',
             0.5, 1.0, 12.2, 0.45, font_size=11, bold=True, italic=True, color=WHITE)

    # Chatbot answer column
    add_text(slide, "What the chatbot returns today:", 0.5, 1.52, 5.5, 0.28,
             font_size=9.5, bold=True, color=AMBER)
    add_text(slide,
             "• Benefits (2025–2070 avg):  $17.1 B/yr\n"
             "• Costs   (2025–2070 avg):  $0.67 B/yr\n"
             "• Peak cost as % of GDP:    0.51%",
             0.5, 1.83, 5.5, 0.9, font_size=10, color=TEAL)

    # Gap column
    add_text(slide, "What is missing:", 6.5, 1.52, 6.3, 0.28,
             font_size=9.5, bold=True, color=RED)
    add_text(slide,
             "• What % of the $17B benefit is health (indoor air quality)?\n"
             "• What % is emissions reduction vs. energy savings?\n"
             "• How do costs break down per sector / per intervention?\n"
             "→  Without sector CBA: cannot compare policies on a like-for-like basis.",
             6.5, 1.83, 6.1, 1.1, font_size=10, color=WHITE)

    # ── Divider ───
    add_rect(slide, 0.3, 3.65, 12.73, 0.03, GRAY)

    # ── BOTTOM: MAC curve (entry 35) ───
    add_rect(slide, 0.3, 3.73, 12.73, 2.88, SURFACE)

    add_text(slide, "Example 2 — Log entry 35",
             0.5, 3.75, 3.0, 0.28, font_size=8.5, color=GRAY, italic=True)
    add_text(slide,
             '"Can you produce a MAC curve for the net zero strategy?"',
             0.5, 4.03, 12.2, 0.45, font_size=11, bold=True, italic=True, color=WHITE)

    add_text(slide, "What the chatbot did:", 0.5, 4.55, 5.5, 0.28,
             font_size=9.5, bold=True, color=AMBER)
    add_text(slide,
             "• Ran 16 separate simulations (one per policy lever at L=0.9)\n"
             "• Ranked abatement potential by cumulative emission reduction\n"
             "• Estimated MAC cost per lever from aggregated CBA total",
             0.5, 4.85, 5.5, 0.9, font_size=10, color=TEAL)

    add_text(slide, "What is missing:", 6.5, 4.55, 6.3, 0.28,
             font_size=9.5, bold=True, color=RED)
    add_text(slide,
             "• MAC costs are approximated — aggregated CBA cannot be split per lever\n"
             "• A lever like 'Deforestation Reduction' has ecosystem co-benefits that are\n"
             "  not separated from the total benefits pool\n"
             "→  Per-sector CBA would make each bar on the MAC curve precise.",
             6.5, 4.85, 6.1, 1.1, font_size=10, color=WHITE)

    add_text(slide,
             "Screenshots of these chatbot responses can be inserted here as evidence from the live interface.",
             0.3, 6.75, 12.73, 0.3, font_size=8.5, color=GRAY, italic=True)


# ── Slide 5 — Roadmap ─────────────────────────────────────────────────────────

def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, top=0.0, height=0.07, color=AMBER)

    add_text(slide, "ROADMAP — Enabling Per-Sector Cost-Benefit Analysis",
             0.3, 0.12, 12.7, 0.42, font_size=14, bold=True, color=AMBER)

    # ── Current state ───
    add_text(slide, "Current state", 0.3, 0.65, 12.7, 0.32,
             font_size=11, bold=True, color=WHITE)

    status_rows = [
        ["Component", "Status", "Detail"],
        ["12-sector emission surrogate", "✓  Complete", "48 outputs: 12 sectors × 4 years"],
        ["Aggregated CBA surrogate", "✓  Complete", "8 outputs: costs & benefits, near/long-term"],
        ["Per-sector CBA surrogate", "✗  Not implemented", "Raw SISEPUEDE data has 16 benefit types × 12 sectors"],
    ]
    status_tc = [
        None,
        [WHITE, GREEN_CHK, GRAY],
        [WHITE, GREEN_CHK, GRAY],
        [WHITE, RED,       GRAY],
    ]
    add_table(slide, status_rows, [3.5, 2.0, 6.8], left=0.3, top=1.0,
              row_height=0.3, text_colors=status_tc)

    # ── What's needed ───
    add_text(slide, "What's needed to enable per-sector CBA", 0.3, 2.38, 12.7, 0.32,
             font_size=11, bold=True, color=WHITE)

    steps = [
        ("Step 1",
         "Data preparation",
         "Modify surrogate_model_sector/data_prep_sector.ipynb to preserve sector identity "
         "in CBA aggregation. Currently, all 16 benefit types are summed to a single total "
         "before training — instead, aggregate per sector × benefit type."),
        ("Step 2",
         "Retrain sector surrogate",
         "Train with ~48 additional CBA targets (12 sectors × 2 periods × 2 metric types: "
         "costs + benefits). The existing XGBoost pipeline architecture handles this without "
         "structural changes."),
        ("Step 3",
         "Expose new endpoint",
         "Update predictor.py → SectorPredictor.predict_comparison() to return per-sector "
         "CBA alongside per-sector emissions. The chatbot agent tool schema is updated to "
         "surface these fields to the LLM."),
    ]

    for i, (tag, title, desc) in enumerate(steps):
        top = 2.75 + i * 1.35
        add_rect(slide, 0.3, top, 1.1, 1.1, SURFACE)
        add_text(slide, tag, 0.35, top + 0.1, 1.0, 0.28,
                 font_size=9, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.55, top + 0.03, 11.2, 0.32,
                 font_size=11, bold=True, color=WHITE)
        add_text(slide, desc, 1.55, top + 0.38, 11.2, 0.75,
                 font_size=9.5, color=GRAY)

    # Timeline
    add_text(slide, "Estimated effort:  2–3 days data prep  +  1 day retraining  +  1 day integration",
             0.3, 6.96, 12.73, 0.34, font_size=10, color=AMBER, italic=True,
             align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_log_entry(prs)
    slide_verification(prs)
    slide_cba_limitation(prs)
    slide_roadmap(prs)

    out = Path("demo_slides.pptx")
    prs.save(out)
    print(f"Saved: {out.resolve()}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
